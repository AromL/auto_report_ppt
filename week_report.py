# -*- coding: utf-8 -*-
# 数据、基础部分
# 连接数据库
from sqlalchemy import create_engine as ce
# 日期时间函数（用来获取当前时间）
import datetime as dt
# 读取、处理数据表
import pandas as pd
# 文件复制
import shutil
# 数学处理函数
import math
# JSON字符串读取和生成
import json
# 时间函数
import time
# 文件目录遍历
import os
# ppt部分
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData


def get_engine():
    with open('./needed/database.txt', 'r') as f:  # r只读 w只写 a追加 rw读写 b二进制操作
        db = json.loads(f.read())
    engine = ce('mysql+mysqlconnector://'+db['usr']+':'+db['pwd']+'@'+db['host']+'/'+db['db'])
    # mysql+mysqlconnector://{usr}:{pwd}@{host}:{port}/{db}
    return engine


def run_sql_query():
    dt7 = (dt.date.today() - dt.timedelta(days=dt.date.today().isocalendar()[2])).strftime('%Y-%m-%d')
    end_date = dt.datetime.strptime(dt7, '%Y-%m-%d')
    dt6 = (end_date+dt.timedelta(days=-1)).strftime('%Y-%m-%d')
    dt5 = (end_date+dt.timedelta(days=-2)).strftime('%Y-%m-%d')
    dt4 = (end_date+dt.timedelta(days=-3)).strftime('%Y-%m-%d')
    dt3 = (end_date+dt.timedelta(days=-4)).strftime('%Y-%m-%d')
    dt2 = (end_date+dt.timedelta(days=-5)).strftime('%Y-%m-%d')
    dt1 = (end_date+dt.timedelta(days=-6)).strftime('%Y-%m-%d')

    date_list = list()
    year_str = dt1.split('-')[0]
    first_day_of_year = dt.datetime.strptime(year_str+'-01-01', '%Y-%m-%d')
    days_delta = 7-first_day_of_year.isocalendar()[2]
    first_sumday_of_year = first_day_of_year+dt.timedelta(days=days_delta)
    current_date = first_sumday_of_year
    while current_date <= end_date:
        date_list.append("'"+current_date.strftime('%Y-%m-%d')+"'")
        current_date += dt.timedelta(days=7)
    date_str = ','.join(date_list)  # 用于 sql_role_times 中

    roles = ['店长', '店员', '管理']

    sql_rihuo = "select area as '地区', "\
        "round(tran_avg, 2) as '平均日活' "\
    "from "\
        "( "\
        "select a.area, "\
            "dt1, dt2, dt3, dt4, dt5, dt6, dt7, "\
            "dt1_all, dt2_all, dt3_all, dt4_all, dt5_all, dt6_all, dt7_all, "\
            "ROUND(dt1/dt1_all, 2) as tran_dt1, "\
            "ROUND(dt2/dt2_all, 2) as tran_dt2, "\
            "ROUND(dt3/dt3_all, 2) as tran_dt3, "\
            "ROUND(dt4/dt4_all, 2) as tran_dt4, "\
            "ROUND(dt5/dt5_all, 2) as tran_dt5, "\
            "ROUND(dt6/dt6_all, 2) as tran_dt6, "\
            "ROUND(dt7/dt7_all, 2) as tran_dt7, "\
            "(dt1/dt1_all+dt2/dt2_all+dt3/dt3_all+dt4/dt4_all+dt5/dt5_all+dt6/dt6_all+dt7/dt7_all)/7 as tran_avg "\
        "from "\
            "( "\
                "SELECT area, "\
                    "sum(if(dt = '"+dt1+"',day_user_count,0)) AS dt1, "\
                    "sum(if(dt = '"+dt2+"',day_user_count,0)) AS dt2, "\
                    "sum(if(dt = '"+dt3+"',day_user_count,0)) AS dt3, "\
                    "sum(if(dt = '"+dt4+"',day_user_count,0)) AS dt4, "\
                    "sum(if(dt = '"+dt5+"',day_user_count,0)) AS dt5, "\
                    "sum(if(dt = '"+dt6+"',day_user_count,0)) AS dt6, "\
                    "sum(if(dt = '"+dt7+"',day_user_count,0)) AS dt7 "\
                "from ystp_app_day_user "\
                "where dt between '"+dt1+"' and '"+dt7+"' and area not in ('华南','华东') "\
                    "and day_app_name ='all_user_count' "\
                "group by area "\
                "union all "\
                "SELECT '总计' area, "\
                    "sum(if(dt = '"+dt1+"',day_user_count,0)) AS dt1, "\
                    "sum(if(dt = '"+dt2+"',day_user_count,0)) AS dt2, "\
                    "sum(if(dt = '"+dt3+"',day_user_count,0)) AS dt3, "\
                    "sum(if(dt = '"+dt4+"',day_user_count,0)) AS dt4, "\
                    "sum(if(dt = '"+dt5+"',day_user_count,0)) AS dt5, "\
                    "sum(if(dt = '"+dt6+"',day_user_count,0)) AS dt6, "\
                    "sum(if(dt = '"+dt7+"',day_user_count,0)) AS dt7 "\
                "from ystp_app_day_user "\
                "where dt between '"+dt1+"' and '"+dt7+"' and area not in ('华南','华东') "\
                    "and day_app_name ='all_user_count' "\
            ") a "\
        "left join "\
            "( "\
                "select area, "\
                    "sum(if(dt = '"+dt1+"',all_user_count,0)) AS dt1_all, "\
                    "sum(if(dt = '"+dt2+"',all_user_count,0)) AS dt2_all, "\
                    "sum(if(dt = '"+dt3+"',all_user_count,0)) AS dt3_all, "\
                    "sum(if(dt = '"+dt4+"',all_user_count,0)) AS dt4_all, "\
                    "sum(if(dt = '"+dt5+"',all_user_count,0)) AS dt5_all, "\
                    "sum(if(dt = '"+dt6+"',all_user_count,0)) AS dt6_all, "\
                    "sum(if(dt = '"+dt7+"',all_user_count,0)) AS dt7_all "\
                "from ystp_allapp_all_user "\
                "where dt between '"+dt1+"' and '"+dt7+"' and area not in ('华南','华东') "\
                "group by area "\
                "union all "\
                "select '总计' area, "\
                    "sum(if(dt = '"+dt1+"',all_user_count,0)) AS dt1_all, "\
                    "sum(if(dt = '"+dt2+"',all_user_count,0)) AS dt2_all, "\
                    "sum(if(dt = '"+dt3+"',all_user_count,0)) AS dt3_all, "\
                    "sum(if(dt = '"+dt4+"',all_user_count,0)) AS dt4_all, "\
                    "sum(if(dt = '"+dt5+"',all_user_count,0)) AS dt5_all, "\
                    "sum(if(dt = '"+dt6+"',all_user_count,0)) AS dt6_all, "\
                    "sum(if(dt = '"+dt7+"',all_user_count,0)) AS dt7_all "\
                "from ystp_allapp_all_user "\
                "where dt between '"+dt1+"' and '"+dt7+"' and area not in ('华南','华东') "\
            ") b on a.area=b.area) t"

    sql_rihuo_tu = "select a.dt as '日期', "\
        "all_user_count as '日安装', "\
        "day_user_count as '日活跃', "\
        "ROUND(day_user_count/all_user_count, 2) as '日活率' "\
    "from "\
        "( "\
            "SELECT dt, sum(day_user_count) as day_user_count "\
            "from ystp_app_day_user "\
            "where dt between '"+dt1+"' and '"+dt7+"' and area not in ('华南','华东') "\
                "and day_app_name ='all_user_count' "\
            "group by dt "\
        ") a "\
        "left join "\
        "( "\
            "select dt, sum(all_user_count) as all_user_count "\
            "from ystp_allapp_all_user "\
            "where dt between '"+dt1+"' and '"+dt7+"' and area not in ('华南','华东') "\
            "group by dt "\
        ") b on a.dt=b.dt"

    sql_week_data = "select a.area as '地区', "\
        "all_user_count as '周安装', "\
        "new_user_count as '新增安装', "\
        "week_user_count as '周活跃', "\
        "ROUND(week_user_count/all_user_count, 2) as '周活跃率' "\
    "from "\
        "( "\
            "select area, "\
                "sum(week_user_count) as week_user_count "\
            "from ystp_app_week_user "\
            "where dt='"+dt7+"' and area not in ('华南','华东') "\
                "and week_app_name='all_user_count' "\
            "group by area "\
            "union all "\
            "select '总计' area, "\
                "sum(week_user_count) as week_user_count "\
            "from ystp_app_week_user "\
            "where dt='"+dt7+"' and area not in ('华南','华东') "\
                "and week_app_name='all_user_count' "\
        ") a "\
        "left join "\
        "( "\
            "select area, "\
                "sum(all_user_count) as all_user_count, "\
                "sum(new_user_count) as new_user_count "\
            "from ystp_allapp_all_user "\
            "where dt='"+dt7+"' and area not in ('华南','华东') "\
            "group by area "\
            "union all "\
            "select '总计' area, "\
                "sum(all_user_count) as all_user_count, "\
                "sum(new_user_count) as new_user_count "\
            "from ystp_allapp_all_user "\
            "where dt='"+dt7+"' and area not in ('华南','华东') "\
        ") b on a.area=b.area"

    sql_week_avg_list = list()
    sql_role_week_avg = "select dayofweek(t1.dt) '周', "\
        "round(t1.rel1/t2.rel2, 2) rate "\
    "from "\
        "( "\
            "select dt, sum(day_user_count) rel1 "\
            "from ystp_app_day_user "\
            "where day_app_name = 'all_user_count' and flag2 = '{}' and dt >= '"+dt1+"' "\
            "group by dt "\
        ") t1, "\
        "( "\
            "select dt, sum(all_user_count) rel2 "\
            "from ystp_allapp_all_user "\
            "where flag2 = '{}' and dt >= '"+dt1+"' "\
            "group by dt "\
        ") t2 "\
    "where t1.dt = t2.dt "\
    "group by 周"

    for role in roles:
        now_sql = sql_role_week_avg.format(role, role)
        sql_week_avg_list.append(now_sql)

    sql_rihuo_list = list()
    sql_role_rihuo = "select area as '地区', "\
        "round(tran_avg, 2) as '平均日活' "\
    "from "\
        "( "\
            "select a.area, "\
                "dt1, dt2, dt3, dt4, dt5, dt6, dt7, "\
                "dt1_all, dt2_all, dt3_all, dt4_all, dt5_all, dt6_all, dt7_all, "\
                "ROUND(dt1/dt1_all, 2) as tran_dt1, "\
                "ROUND(dt2/dt2_all, 2) as tran_dt2, "\
                "ROUND(dt3/dt3_all, 2) as tran_dt3, "\
                "ROUND(dt4/dt4_all, 2) as tran_dt4, "\
                "ROUND(dt5/dt5_all, 2) as tran_dt5, "\
                "ROUND(dt6/dt6_all, 2) as tran_dt6, "\
                "ROUND(dt7/dt7_all, 2) as tran_dt7, "\
                "(dt1/dt1_all+dt2/dt2_all++dt3/dt3_all+dt4/dt4_all+dt5/dt5_all+dt6/dt6_all+dt7/dt7_all)/7 as tran_avg "\
            "from "\
            "( "\
                "SELECT area, "\
                    "sum(if(dt = '"+dt1+"',day_user_count,0)) AS dt1, "\
                    "sum(if(dt = '"+dt2+"',day_user_count,0)) AS dt2, "\
                    "sum(if(dt = '"+dt3+"',day_user_count,0)) AS dt3, "\
                    "sum(if(dt = '"+dt4+"',day_user_count,0)) AS dt4, "\
                    "sum(if(dt = '"+dt5+"',day_user_count,0)) AS dt5, "\
                    "sum(if(dt = '"+dt6+"',day_user_count,0)) AS dt6, "\
                    "sum(if(dt = '"+dt7+"',day_user_count,0)) AS dt7 "\
                "from ystp_app_day_user "\
                "where dt between '"+dt1+"' and '"+dt7+"' and area not in ('华南','华东') "\
                    "and flag2 = '{}' "\
                    "and day_app_name ='all_user_count' "\
                "group by area "\
                "union all "\
                "SELECT '总计' area, "\
                    "sum(if(dt = '"+dt1+"',day_user_count,0)) AS dt1, "\
                    "sum(if(dt = '"+dt2+"',day_user_count,0)) AS dt2, "\
                    "sum(if(dt = '"+dt3+"',day_user_count,0)) AS dt3, "\
                    "sum(if(dt = '"+dt4+"',day_user_count,0)) AS dt4, "\
                    "sum(if(dt = '"+dt5+"',day_user_count,0)) AS dt5, "\
                    "sum(if(dt = '"+dt6+"',day_user_count,0)) AS dt6, "\
                    "sum(if(dt = '"+dt7+"',day_user_count,0)) AS dt7 "\
                "from ystp_app_day_user "\
                "where dt between '"+dt1+"' and '"+dt7+"' and area not in ('华南','华东') "\
                    "and flag2 = '{}' "\
                    "and day_app_name ='all_user_count' "\
            ") a "\
            "left join "\
            "( "\
                "select area, "\
                    "sum(if(dt = '"+dt1+"',all_user_count,0)) AS dt1_all, "\
                    "sum(if(dt = '"+dt2+"',all_user_count,0)) AS dt2_all, "\
                    "sum(if(dt = '"+dt3+"',all_user_count,0)) AS dt3_all, "\
                    "sum(if(dt = '"+dt4+"',all_user_count,0)) AS dt4_all, "\
                    "sum(if(dt = '"+dt5+"',all_user_count,0)) AS dt5_all, "\
                    "sum(if(dt = '"+dt6+"',all_user_count,0)) AS dt6_all, "\
                    "sum(if(dt = '"+dt7+"',all_user_count,0)) AS dt7_all "\
                "from ystp_allapp_all_user "\
                "where dt between '"+dt1+"' and '"+dt7+"' and area not in ('华南','华东') "\
                    "and flag2 = '{}' "\
                "group by area "\
                "union all "\
                "select '总计' area, "\
                    "sum(if(dt = '"+dt1+"',all_user_count,0)) AS dt1_all, "\
                    "sum(if(dt = '"+dt2+"',all_user_count,0)) AS dt2_all, "\
                    "sum(if(dt = '"+dt3+"',all_user_count,0)) AS dt3_all, "\
                    "sum(if(dt = '"+dt4+"',all_user_count,0)) AS dt4_all, "\
                    "sum(if(dt = '"+dt5+"',all_user_count,0)) AS dt5_all, "\
                    "sum(if(dt = '"+dt6+"',all_user_count,0)) AS dt6_all, "\
                    "sum(if(dt = '"+dt7+"',all_user_count,0)) AS dt7_all "\
                "from ystp_allapp_all_user "\
                "where dt between '"+dt1+"' and '"+dt7+"' and area not in ('华南','华东') "\
                    "and flag2 = '{}' "\
            ") b on a.area=b.area "\
        ") t"

    for role in roles:
        now_sql = sql_role_rihuo.format(role, role, role, role)
        sql_rihuo_list.append(now_sql)

    sql_rihuo_tu_list = list()
    sql_role_rihuo_tu = "select a.dt as '日期', "\
        "all_user_count as '日安装', "\
        "day_user_count as '日活跃', "\
        "ROUND(day_user_count/all_user_count, 2) as '日活率' "\
    "from "\
        "( "\
            "SELECT dt, sum(day_user_count) as day_user_count "\
            "from ystp_app_day_user "\
            "where dt between '"+dt1+"' and '"+dt7+"' and area not in ('华南','华东') "\
                "and flag2 = '{}' "\
                "and day_app_name ='all_user_count' "\
            "group by dt "\
        ") a "\
        "left join "\
        "( "\
            "select dt, sum(all_user_count) as all_user_count "\
            "from ystp_allapp_all_user "\
            "where dt between '"+dt1+"' and '"+dt7+"' and area not in ('华南','华东') "\
                "and flag2 = '{}' "\
            "group by dt "\
        ") b on a.dt=b.dt"

    for role in roles:
        now_sql = sql_role_rihuo_tu.format(role, role)
        sql_rihuo_tu_list.append(now_sql)

    sql_week_data_list = list()
    sql_role_week_data = "select a.area as '地区', "\
        "all_user_count as '周安装', "\
        "new_user_count as '新增安装', "\
        "week_user_count as '周活跃', "\
        "ROUND(week_user_count/all_user_count, 2) as '周活跃率' "\
    "from "\
        "( "\
            "select area, "\
                "sum(week_user_count) as week_user_count "\
            "from ystp_app_week_user "\
            "where dt='"+dt7+"' and area not in ('华南','华东') "\
                "and flag2 = '{}' "\
                "and week_app_name='all_user_count' "\
            "group by area "\
            "union all "\
            "select '总计' area, "\
                "sum(week_user_count) as week_user_count "\
            "from ystp_app_week_user "\
            "where dt='"+dt7+"' and area not in ('华南','华东') "\
                "and flag2 = '{}' "\
                "and week_app_name='all_user_count' "\
        ") a "\
        "left join "\
        "( "\
            "select area, "\
                "sum(all_user_count) as all_user_count, "\
                "sum(new_user_count) as new_user_count "\
            "from ystp_allapp_all_user "\
            "where dt='"+dt7+"' and area not in ('华南','华东') "\
                "and flag2 = '{}' "\
            "group by area "\
            "union all "\
            "select '总计' area, "\
                "sum(all_user_count) as all_user_count, "\
                "sum(new_user_count) as new_user_count "\
            "from ystp_allapp_all_user "\
            "where dt='"+dt7+"' and area not in ('华南','华东') "\
                "and flag2 = '{}' "\
        ") b on a.area=b.area"

    for role in roles:
        now_sql = sql_role_week_data.format(role, role, role, role)
        sql_week_data_list.append(now_sql)

    except_keywords_list = [
        '店员管理', '百丽统一登录(生产勿动)', '订单管理', '核销管理', '区域筛选', '考勤打卡',
        '移动收银', '内容平台', '培训君', '内容管理', '门店对账', '考试君', '考勤看板', '订单管理',
        '核销管理', '通讯录管理', '用户授权', '打卡', '学霸君', '设备解绑', '轻应用示例'
    ]
    except_keywords_list = ["'"+item+"'" for item in except_keywords_list]
    except_keywords_str = ', '.join(except_keywords_list)  # 用于 sql_role_apps_rank 中

    sql_apps_rank_list = list()
    sql_role_apps_rank = "select week_app_name as '轻应用', "\
        "ROUND(area1, 2) as '沪浙', "\
        "ROUND(area2, 2) as '华北一区', "\
        "ROUND(area3, 2) as '华北二区', "\
        "ROUND(area4, 2) as '华南一区', "\
        "ROUND(area5, 2) as '华南二区', "\
        "ROUND(area6, 2) as '华中', "\
        "ROUND(area7, 2) as '鲁豫东北', "\
        "ROUND(area8, 2) as '苏皖', "\
        "ROUND(area9, 2) as '西北', "\
        "ROUND(area10, 2) as '西南', "\
        "ROUND(area11, 2) as '云贵', "\
        "ROUND(area12, 2) as '体总', "\
        "ROUND(tran_avg, 2) as '平均值' "\
    "from "\
        "( "\
            "select week_app_name, "\
                "area1, area2, area3, area4, area5, area6, area7, area8, area9, area10, area11, area12, "\
                "(area1+area2+area3+area4+area5+area6+area7+area8+area9+area10+area11+area12)/12 tran_avg "\
            "from "\
            "( "\
                "SELECT week_app_name, "\
                    "sum(if(t0.area = '沪浙',app_count/all_count,0)) AS area1, "\
                    "sum(if(t0.area = '华北一区',app_count/all_count,0)) AS area2, "\
                    "sum(if(t0.area = '华北二区',app_count/all_count,0)) AS area3, "\
                    "sum(if(t0.area = '华南一区',app_count/all_count,0)) AS area4, "\
                    "sum(if(t0.area = '华南二区',app_count/all_count,0)) AS area5, "\
                    "sum(if(t0.area = '华中',app_count/all_count,0)) AS area6, "\
                    "sum(if(t0.area = '鲁豫',app_count/all_count,0)) AS area7, "\
                    "sum(if(t0.area = '苏皖',app_count/all_count,0)) AS area8, "\
                    "sum(if(t0.area = '西北',app_count/all_count,0)) AS area9, "\
                    "sum(if(t0.area = '西南',app_count/all_count,0)) AS area10, "\
                    "sum(if(t0.area = '云贵',app_count/all_count,0)) AS area11, "\
                    "sum(if(t0.area = '体总',app_count/all_count,0)) AS area12 "\
                "from "\
                "( "\
                    "select week_app_name, "\
                        "area, "\
                        "sum(week_user_count) as app_count "\
                    "from ystp_app_week_user "\
                    "where dt = '"+dt7+"' "\
                        "and flag2 = '{}' and week_app_name !='' "\
                        "and week_app_name not in ('all_user_count','all_data_count','from_jpush') "\
                        "and week_app_name not like '%测试%' and week_app_name not like '%开发%' "\
                        "and week_app_name not like 'HR%' and week_app_name not like '天眼%' "\
                        "and week_app_name not like '%仪表盘' and week_app_name not like '小程序%' "\
                        "and week_app_name not like '我的%' "\
                        "and week_app_name not in("+except_keywords_str+") "\
                    "group by week_app_name, area "\
                ") t0 "\
                "left join "\
                "( "\
                    "select area, "\
                        "sum(all_user_count) as all_count "\
                    "from ystp_allapp_all_user  "\
                    "where dt = '"+dt7+"' "\
                        "and flag2 = '{}' "\
                    "group by area "\
                ") t1 on t0.area=t1.area "\
                "group by week_app_name "\
            ") t "\
            "order by tran_avg desc "\
        ") t"

    for role in roles:
        now_sql = sql_role_apps_rank.format(role, role, role, role)
        sql_apps_rank_list.append(now_sql)

    week_num = year_str+'W'+str(end_date.isocalendar()[1])  # 用于 sql_role_times 中
    sql_times_list = list()
    sql_role_times = "SELECT days_count AS '天数', "\
	    "sum(IF(period_sdate = '"+week_num+"',human_count,0)) AS "+week_num+" "\
    "FROM "\
        "( "\
            "SELECT period_sdate, "\
                "days_count, "\
                "human_count "\
            "FROM "\
                "( "\
                    "SELECT period_sdate, "\
                        "days_count, "\
                        "sum(human_count) AS human_count "\
                    "FROM ystp_times_user "\
                    "WHERE all_count = 'all_user_count' "\
                    "AND flag2 = '{}' "\
                    "GROUP BY period_sdate, days_count "\
                    "UNION ALL "\
                        "SELECT a.period_sdate, "\
                            "'0' AS days_count, "\
                            "all_user_count - human_all AS human_count "\
                        "FROM "\
                            "( "\
                                "SELECT period_sdate, "\
                                    "days_count, "\
                                    "sum(human_count) AS human_count "\
                                "FROM ystp_times_user "\
                                "WHERE all_count = 'all_user_count' "\
                                "AND flag2 = '{}' "\
                                "GROUP BY period_sdate, days_count "\
                            ") a "\
                        "LEFT JOIN ( "\
                            "SELECT period_sdate, "\
                                "sum(human_count) AS human_all "\
                            "FROM "\
                                "( "\
                                    "SELECT period_sdate, "\
                                        "days_count, "\
                                        "sum(human_count) AS human_count "\
                                    "FROM ystp_times_user "\
                                    "WHERE all_count = 'all_user_count' "\
                                    "AND flag2 = '{}' "\
                                    "GROUP BY period_sdate, days_count "\
                                ") tmp "\
                            "GROUP BY period_sdate "\
                        ") b ON a.period_sdate = b.period_sdate "\
                        "LEFT JOIN ( "\
                            "SELECT dt, "\
                                "period_sdate, "\
                                "sum(all_user_count) AS all_user_count "\
                            "FROM ystp_allapp_all_user "\
                            "WHERE dt <> '"+dt7+"' "\
                            "AND flag2 = '{}' "\
                            "GROUP BY dt "\
                        ") c ON a.period_sdate = c.period_sdate "\
                        "GROUP BY "\
                            "period_sdate "\
                ") t "\
            "ORDER BY period_sdate, days_count "\
        ") t "\
    "GROUP BY days_count "\
    "ORDER BY days_count desc"

    for role in roles:
        now_sql = sql_role_times.format(role, role, role, role)
        sql_times_list.append(now_sql)

    excel_data_dict = {
        '日活数据': [1, sql_rihuo, 0, 0],  # page 2
        '周安装-新增-活跃': [1, sql_week_data, 0, 15],  # page 2
        '日活-店长': [2, sql_rihuo_list[0], 0, 31],  # 在第 29 行添加部分名称 page 3
        '周安装-新增-活跃-店长': [2, sql_week_data_list[0], 0, 47],  # 在第 44 行添加部分名称 page 3
        '应用排行百分比-大区-店长': [3, sql_apps_rank_list[0], 0, 63],  # 在第 59 行添加部分名称 page 4
        '日活-店员': [5, sql_rihuo_list[1], 3, 31],  # page 6
        '周安装-新增-活跃-店员': [5, sql_week_data_list[1], 6, 47],  # page 6
        '应用排行百分比-大区-店员': [6, sql_apps_rank_list[1], 0, 79],  # 在第 74 行添加部分名称 page 7
        '日活-管理': [8, sql_rihuo_list[2], 6, 31],  # page 9
        '周安装-新增-活跃-管理': [8, sql_week_data_list[2], 12, 47],  # page 9
        '应用排行百分比-大区-管理': [9, sql_apps_rank_list[2], 0, 95],  # 在第 89 行添加部分名称 page 10
    }
    chart_data_dict = {
        '日活-图': [0, sql_rihuo_tu],  # page 2
        '日活-图-店长': [0, sql_rihuo_tu_list[0]],  # page 3
        '日活-图-店员': [0, sql_rihuo_tu_list[1]],  # page 6
        '日活-图-管理': [0, sql_rihuo_tu_list[2]],  # page 9
        '使用频次-店长': [1, sql_times_list[0]],  # page 5
        '使用频次-店员': [1, sql_times_list[1]],  # page 8
        '使用频次-管理': [1, sql_times_list[2]]  # page 11
    }

    engine = get_engine()
    table_data = dict()
    today_date_str = dt.date.today().strftime('%Y-%m-%d')
    # 写入 周报数据.xls
    folder_path = './outputs/'+today_date_str
    if not os.path.exists(folder_path):
        os.makedirs(folder_path)
    file_dir = folder_path+'/周报.xls'
    writer1 = pd.ExcelWriter(file_dir)
    for name, items in excel_data_dict.items():
        data = pd.read_sql_query(items[1], engine)
        if len(items) == 4:
            title = name.split('-')[-1]
            if title in roles:
                name_val = pd.DataFrame([title])
                name_val.to_excel(writer1, index=False, header=False, startcol=items[2], startrow=items[3] - 1)
            data.to_excel(writer1, index=False, startcol=items[2], startrow=items[3])
        table_data[name] = data  # 格式为 页码、数据
    writer1.save()
    # 写入 history_data.xls
    chart_data = dict()
    history_output_dir = './outputs/'+today_date_str+'/history_data.xls'
    writer2 = pd.ExcelWriter(history_output_dir)
    for name, items in chart_data_dict.items():
        data = pd.read_sql_query(items[1], engine)
        history_data = pd.read_excel('./needed/history_base.xls', sheet_name=name)
        # 如果新的到的数据在历史数据中不存在，则追加到历史数据
        if items[0] == 0:  # 第一种历史数据 -> 日活数据
            history_date = dt.datetime.strptime(history_data.iloc[-1, 0], '%Y/%m/%d')
            data_date = dt.datetime.strptime(data.iloc[-1, 0], '%Y-%m-%d')
            if history_date != data_date:
                data['星期&日活'] = None
                data['星期'] = ['星期一', '星期二', '星期三', '星期四', '星期五', '星期六', '星期日']
                new_data = pd.concat([history_data, data], ignore_index=True, sort=False)
                new_data['日期'] = pd.to_datetime(new_data['日期'])
                case1 = new_data['日期'] < dt.datetime.strptime('2019-10-7', '%Y-%m-%d')
                case2 = new_data['日期'] > dt.datetime.strptime('2019-10-13', '%Y-%m-%d')
                week_avg = new_data[case1 | case2].groupby('星期')['日活率'].mean()
                last_7 = new_data.tail(7)
                for now_row in last_7[['星期', '星期&日活']].iterrows():
                    # row[0] -> index, row[1]['星期'] -> 星期几
                    new_data.iloc[now_row[0], -2] = '%.2f' % (week_avg[now_row[1]['星期']])
                new_data['星期&日活'] = pd.to_numeric(new_data['星期&日活'])
                new_data['日期'] = new_data['日期'].dt.strftime('%Y/%m/%d')
                history_data = new_data
            history_data.to_excel(writer2, sheet_name=name, index=False)
            chart_data[name] = history_data.tail(31)
        elif items[0] == 1:  # 第二种历史数据 -> 使用频次数据
            if week_num not in history_data.keys():
                history_data[week_num] = data[week_num]
            history_data.to_excel(writer2, sheet_name=name, index=False)
            chart_data[name] = history_data.iloc[:, [0, -11, -10, -9, -8, -7, -6, -5, -4, -3, -2, -1]]
    writer2.save()

    return table_data, chart_data


def run_save_ppt(table, chart):
    today_date = dt.date.today()
    sunday = dt.date.today() - dt.timedelta(days=dt.date.today().isocalendar()[2])
    monday = (sunday+dt.timedelta(days=-6)).strftime('%Y.%m.%d')
    sunday = sunday.strftime('%Y.%m.%d')
    # 根据模板生成 ppt
    prs = Presentation('./needed/template.pptx')
    slides = prs.slides

    # 第一页
    slide1 = slides[0]
    shapes1 = slide1.shapes
    shapes1[1].text_frame.paragraphs[0].clear()  # 清除原格式
    set_sm_title1 = shapes1[1].text_frame.paragraphs[0].add_run()  # 重新添加格式
    # set_sm_title1.text = dt1+' - '+dt7
    set_sm_title1.text = monday+' - '+sunday
    set_sm_title1.font.name = '华文楷体'
    set_sm_title1.font.size = Pt(20)
    set_sm_title1.font.bold = True
    set_sm_title1.font.color.rgb = RGBColor(0x95, 0x37, 0x35)

    # 第二页、第三页、第六页、第九页 -> table、chart -> 都有两个数据表 -> 构建数据结构
    table2_arr = [
        ['周安装-新增-活跃', '周安装-新增-活跃-店长', '周安装-新增-活跃-店员', '周安装-新增-活跃-管理'],
        ['日活数据', '日活-店长', '日活-店员', '日活-管理']
    ]
    chart2_arr = ['日活-图', '日活-图-店长', '日活-图-店员', '日活-图-管理']
    pages = [1, 2, 5, 8]
    for i in range(len(pages)):
        shapes2 = slides[pages[i]].shapes
        shape_chart, shape_table = None, None
        for shape in shapes2:
            if '图表' in shape.name:
                shape_chart = shape
            if '表格' in shape.name:
                shape_table = shape
        # # 左侧表格
        now_table = table[table2_arr[0][i]]
        for item in now_table.iterrows():  # 前5列数据
            for j in range(len(item[1])):
                cell = shape_table.table.cell(item[0]+1, j)
                cell.margin_top = cell.margin_bottom = Inches(.06)
                cell.text_frame.paragraphs[0].clear()  # 清除原格式
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                set_cell = cell.text_frame.paragraphs[0].add_run()
                if j == 0:
                    ins_str = item[1][j]
                elif j < 4:
                    ins_str = str(math.floor(item[1][j]))
                else:
                    ins_str = str(math.floor(item[1][j]*100))+'%'
                set_cell.text = ins_str
                set_cell.font.name = '华文楷体'
                set_cell.font.size = Pt(8)
        now_table = table[table2_arr[1][i]]['平均日活']
        for ind, item in enumerate(now_table):  # 最后1列数据
            cell = shape_table.table.cell(ind + 1, 5)
            cell.margin_top = cell.margin_bottom = Inches(.06)
            cell.text_frame.paragraphs[0].clear()  # 清除原格式
            set_cell = cell.text_frame.paragraphs[0].add_run()
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            set_cell.text = str(math.floor(item*100))+'%'
            set_cell.font.name = '华文楷体'
            set_cell.font.size = Pt(8)
        # # 右侧图表
        chart2 = ChartData()
        now_chart = chart[chart2_arr[i]]
        chart2.categories = [item[1] for item in now_chart['日期'].iteritems()]
        chart2.add_series('日安装',    tuple([item[1] for item in now_chart['日安装'].iteritems()]))
        chart2.add_series('日活跃',    tuple([item[1] for item in now_chart['日活跃'].iteritems()]))
        chart2.add_series('日活率',    tuple([item[1] for item in now_chart['日活率'].iteritems()]))
        tmp_tuple = tuple([None if math.isnan(item[1]) else item[1] for item in now_chart['星期&日活'].iteritems()])
        chart2.add_series('星期&日活', tmp_tuple)
        shape_chart.chart.replace_data(chart2)

    table3_arr = ['应用排行百分比-大区-店长', '应用排行百分比-大区-店员', '应用排行百分比-大区-管理']
    # 第四页、第七页、第十页
    pages = [3, 6, 9]
    for i in range(len(pages)):
        shapes3 = slides[pages[i]].shapes
        shape_chart, shape_table = None, None
        for shape in shapes3:
            if '图表' in shape.name:
                shape_chart = shape
            if '表格' in shape.name:
                shape_table = shape
        # # 左侧表格
        now_table = table[table3_arr[i]].drop(columns=['平均值'])
        now_table = now_table.drop(columns=['体总']) if i < 2 else now_table
        now_role = table3_arr[i].split('-')[-1]
        now_table = now_table[0:10] if now_role in ['店长', '管理'] else now_table[0:8]
        for item in now_table.iterrows():  # 前5列数据
            for j in range(len(item[1])):
                cell = shape_table.table.cell(item[0]+1, j)
                cell.margin_top = cell.margin_bottom = Inches(.06)
                cell.text_frame.paragraphs[0].clear()  # 清除原格式
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                set_cell = cell.text_frame.paragraphs[0].add_run()
                ins_str = item[1][j] if type(item[1][j]) == str else str(math.floor(item[1][j]*100))+'%'
                set_cell.text = ins_str
                set_cell.font.name = '华文楷体'
                set_cell.font.size = Pt(8)
        # # 右侧图表
        chart3 = ChartData()
        chart3.categories = [item[1] for item in now_table['轻应用'].iteritems()]
        for key in now_table.keys():
            if key != '轻应用':
                chart3.add_series(key, tuple([float(item[1]) for item in now_table[key].iteritems()]))
        shape_chart.chart.replace_data(chart3)

    chart4_arr = ['使用频次-店长', '使用频次-店员', '使用频次-管理']
    # 第五页、第八页、第十一页
    pages = [4, 7, 10]
    for i in range(len(pages)):
        shapes3 = slides[pages[i]].shapes
        shape_chart = None
        for shape in shapes3:
            if '图表' in shape.name:
                shape_chart = shape
        # # 图表
        now_data = chart[chart4_arr[i]]
        chart4 = ChartData()
        chart4.categories = [key for key in now_data.keys()][1:]
        for row in now_data.iterrows():
            chart4.add_series(str(row[0])+'天', tuple([int(item) for item in row[1]][1:]))
        shape_chart.chart.replace_data(chart4)
    date_str = today_date.strftime('%Y-%m-%d')
    file_name = '大算运营周报'+today_date.strftime('%m.%d')+'（标准）.pptx'
    prs.save('./outputs/'+date_str+'/'+file_name)
    os.remove('./needed/history_base.xls')
    shutil.copy('./outputs/'+date_str+'/history_data.xls', './needed/history_base.xls')


print('--- 开始进行数据读取 ---')
today_date = dt.datetime.today().strftime('%Y-%m-%d')
start_run_time = time.time()
t_data, c_data = run_sql_query()
print('--- 数据读取结束 ---')
print('--- 周报数据已保存到：outputs/'+today_date+' 中 ---')
print('--- 本周历史数据已保存到：outputs/'+today_date+' 中 ---')
end_run_time = time.time()
print('--- 数据处理总计用时：', '%.2f' % (end_run_time-start_run_time), 's ---\n')
print('--- 开始生成并保存PPT ---')
run_save_ppt(t_data, c_data)
print('--- PPT生成结束 ---')
print('--- PPT已保存到：./outputs/'+today_date+' 中 ---')
print('--- 全部历史数据已保存到：./needed/history_base.xls 中 ---')
end_run_time = time.time()
print('--- 全过程总计用时：', '%.2f' % (end_run_time-start_run_time), 's ---')
